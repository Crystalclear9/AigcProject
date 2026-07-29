from __future__ import annotations

import asyncio
import hashlib
import mimetypes
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Awaitable, Callable

from app.schemas.card_refinement import AttachmentDescriptor
from app.services.vivo_ocr import VivoOcrClient

MAX_EXTRACTED_CHARACTERS = 200_000
MAX_PDF_PAGES = 80
MAX_LOCAL_OCR_PAGES = 12
MAX_ARCHIVE_ENTRIES = 10_000
MAX_ARCHIVE_EXPANSION_RATIO = 100


class DocumentExtractionError(RuntimeError):
    def __init__(self, message: str, status: str = "failed") -> None:
        super().__init__(message)
        self.status = status


@dataclass(frozen=True)
class ExtractedDocument:
    descriptor: AttachmentDescriptor
    text: str


OcrHandler = Callable[[bytes], Awaitable[str]]


async def _default_ocr(image_bytes: bytes) -> str:
    lines = await VivoOcrClient().recognize(image_bytes)
    return "\n".join(line.text for line in lines if line.text.strip())


async def extract_document(
    path: Path,
    *,
    name: str,
    declared_mime: str,
    attachment_id: str,
    ocr_handler: OcrHandler | None = None,
) -> ExtractedDocument:
    data = await asyncio.to_thread(path.read_bytes)
    digest = hashlib.sha256(data).hexdigest()
    try:
        detected = detect_document_type(data, name)
    except DocumentExtractionError as error:
        return ExtractedDocument(
            descriptor=AttachmentDescriptor(
                id=attachment_id,
                name=name,
                mime_type="application/octet-stream",
                size_bytes=len(data),
                sha256=digest,
                extraction_status=error.status,
                warning=str(error),
            ),
            text="",
        )
    descriptor = AttachmentDescriptor(
        id=attachment_id,
        name=name,
        mime_type=detected,
        size_bytes=len(data),
        sha256=digest,
    )
    if not mime_matches(declared_mime, detected):
        return ExtractedDocument(
            descriptor=descriptor.model_copy(
                update={
                    "extraction_status": "unsupported",
                    "warning": "文件内容与声明格式不一致",
                }
            ),
            text="",
        )

    try:
        if detected == "application/pdf":
            text, pages, warning = await _extract_pdf(data, ocr_handler or _default_ocr)
        elif detected == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text, pages, warning = await asyncio.to_thread(_extract_docx, path)
        elif detected == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text, pages, warning = await asyncio.to_thread(_extract_pptx, path)
        elif detected == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            text, pages, warning = await asyncio.to_thread(_extract_xlsx, path)
        elif detected in {"text/plain", "text/markdown"}:
            text, pages, warning = _extract_text(data), None, None
        elif detected in {"image/jpeg", "image/png"}:
            text = await (ocr_handler or _default_ocr)(data)
            pages, warning = 1, None
        else:
            raise DocumentExtractionError("不支持的文件格式", "unsupported")
    except DocumentExtractionError as error:
        return ExtractedDocument(
            descriptor=descriptor.model_copy(
                update={
                    "extraction_status": error.status,
                    "warning": str(error),
                }
            ),
            text="",
        )
    except Exception as error:
        return ExtractedDocument(
            descriptor=descriptor.model_copy(
                update={
                    "extraction_status": "failed",
                    "warning": f"解析失败：{type(error).__name__}",
                }
            ),
            text="",
        )

    normalized = normalize_extracted_text(text)
    status = "succeeded" if normalized else "degraded"
    if not normalized and warning is None:
        warning = "未提取到可用文字"
    return ExtractedDocument(
        descriptor=descriptor.model_copy(
            update={
                "extraction_status": status,
                "page_count": pages,
                "extracted_characters": len(normalized),
                "warning": warning,
            }
        ),
        text=normalized,
    )


def detect_document_type(data: bytes, name: str) -> str:
    lowered = name.lower()
    if data.startswith(b"%PDF-"):
        return "application/pdf"
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if data.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if data.startswith(b"PK\x03\x04"):
        return _detect_ooxml(data)
    if b"\x00" not in data[:4096]:
        try:
            data[:65536].decode("utf-8-sig")
        except UnicodeDecodeError:
            pass
        else:
            return "text/markdown" if lowered.endswith((".md", ".markdown")) else "text/plain"
    guessed = mimetypes.guess_type(name)[0]
    raise DocumentExtractionError(
        f"不支持或无法识别的文件格式：{guessed or 'unknown'}",
        "unsupported",
    )


def _detect_ooxml(data: bytes) -> str:
    from io import BytesIO

    try:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            _validate_ooxml_archive(archive)
            names = set(archive.namelist())
    except zipfile.BadZipFile as error:
        raise DocumentExtractionError("Office 文件结构损坏", "unsupported") from error
    if "word/document.xml" in names:
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if "ppt/presentation.xml" in names:
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if "xl/workbook.xml" in names:
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    raise DocumentExtractionError("仅支持 DOCX、PPTX 和 XLSX", "unsupported")


def _validate_ooxml_archive(archive: zipfile.ZipFile) -> None:
    entries = archive.infolist()
    if len(entries) > MAX_ARCHIVE_ENTRIES:
        raise DocumentExtractionError("Office 文件条目过多", "too_large")
    compressed = sum(max(item.compress_size, 1) for item in entries)
    expanded = sum(item.file_size for item in entries)
    if expanded > compressed * MAX_ARCHIVE_EXPANSION_RATIO:
        raise DocumentExtractionError("Office 文件压缩比异常", "too_large")
    if any(
        item.filename.startswith(("/", "\\"))
        or ".." in Path(item.filename).parts
        for item in entries
    ):
        raise DocumentExtractionError("Office 文件包含不安全路径", "unsupported")


def mime_matches(declared: str, detected: str) -> bool:
    normalized = (declared or "").split(";", 1)[0].strip().lower()
    if not normalized or normalized == "application/octet-stream":
        return True
    aliases = {
        "image/jpg": "image/jpeg",
        "text/x-markdown": "text/markdown",
    }
    normalized = aliases.get(normalized, normalized)
    if {normalized, detected} <= {"text/plain", "text/markdown"}:
        return True
    return normalized == detected


async def _extract_pdf(data: bytes, ocr_handler: OcrHandler) -> tuple[str, int, str | None]:
    try:
        import fitz
    except ImportError as error:
        raise DocumentExtractionError("服务端未安装 PDF 解析组件") from error
    try:
        document = fitz.open(stream=data, filetype="pdf")
    except Exception as error:
        raise DocumentExtractionError("PDF 文件损坏或无法打开", "unsupported") from error
    try:
        if document.needs_pass:
            raise DocumentExtractionError("加密 PDF 暂不支持", "password_protected")
        page_count = document.page_count
        if page_count > MAX_PDF_PAGES:
            raise DocumentExtractionError(
                f"PDF 超过 {MAX_PDF_PAGES} 页限制",
                "too_large",
            )
        chunks = [page.get_text("text") for page in document]
        text = "\n".join(chunks)
        warning: str | None = None
        if len(normalize_extracted_text(text)) < max(80, page_count * 20):
            ocr_chunks: list[str] = []
            for index in range(min(page_count, MAX_LOCAL_OCR_PAGES)):
                page = document.load_page(index)
                pixmap = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
                try:
                    ocr_chunks.append(await ocr_handler(pixmap.tobytes("png")))
                except Exception:
                    continue
            if ocr_chunks:
                text = "\n".join(ocr_chunks)
                if page_count > MAX_LOCAL_OCR_PAGES:
                    warning = f"扫描 PDF 仅 OCR 前 {MAX_LOCAL_OCR_PAGES} 页"
            else:
                warning = "扫描 PDF OCR 不可用"
        return text, page_count, warning
    finally:
        document.close()


def _extract_docx(path: Path) -> tuple[str, int | None, str | None]:
    try:
        from docx import Document
    except ImportError as error:
        raise DocumentExtractionError("服务端未安装 DOCX 解析组件") from error
    document = Document(path)
    chunks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            chunks.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(chunks), None, None


def _extract_pptx(path: Path) -> tuple[str, int, str | None]:
    try:
        from pptx import Presentation
    except ImportError as error:
        raise DocumentExtractionError("服务端未安装 PPTX 解析组件") from error
    presentation = Presentation(path)
    chunks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = [
            str(shape.text).strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and str(shape.text).strip()
        ]
        if slide_text:
            chunks.append(f"[第 {index} 页]\n" + "\n".join(slide_text))
    return "\n".join(chunks), len(presentation.slides), None


def _extract_xlsx(path: Path) -> tuple[str, int, str | None]:
    try:
        from openpyxl import load_workbook
    except ImportError as error:
        raise DocumentExtractionError("服务端未安装 XLSX 解析组件") from error
    workbook = load_workbook(path, read_only=True, data_only=True)
    chunks: list[str] = []
    cell_count = 0
    try:
        for sheet in workbook.worksheets[:20]:
            chunks.append(f"[工作表：{sheet.title}]")
            for row in sheet.iter_rows():
                values = [str(cell.value).strip() for cell in row if cell.value is not None]
                if values:
                    chunks.append(" | ".join(values))
                cell_count += len(row)
                if cell_count >= 20_000:
                    return "\n".join(chunks), len(workbook.worksheets), "仅解析前 20000 个单元格"
    finally:
        workbook.close()
    return "\n".join(chunks), len(workbook.worksheets), None


def _extract_text(data: bytes) -> str:
    try:
        return data.decode("utf-8-sig")
    except UnicodeDecodeError as error:
        raise DocumentExtractionError("文本文件必须使用 UTF-8 编码", "unsupported") from error


def normalize_extracted_text(text: str) -> str:
    lines = [" ".join(line.split()) for line in text.replace("\x00", "").splitlines()]
    normalized = "\n".join(line for line in lines if line)
    return normalized[:MAX_EXTRACTED_CHARACTERS]
