from __future__ import annotations

import asyncio
from datetime import datetime, timedelta, timezone
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from app.schemas.card import ActionCard
from app.schemas.card_refinement import CardRefinementPlan, PlanItem, RefinementOptions, UserProfileContext
from app.services.card_refinement_graph import deterministic_plan, validate_plan
from app.services.document_extractor import extract_document


def _card(*, deadline: str | None = None) -> ActionCard:
    return ActionCard(
        id="card-refinement-test",
        card_type="task",
        title="提交课程研究报告",
        summary="完成研究报告并上传学习通",
        deadline=deadline,
        materials=["研究报告 PDF", "数据表"],
        submit_method="学习通",
        source_text="课程通知",
        created_at=datetime.now(timezone.utc),
    )


def test_rule_refinement_schedules_milestones_before_parent_deadline() -> None:
    deadline = datetime.now(timezone.utc) + timedelta(days=5)
    card = _card(deadline=deadline.isoformat())
    plan = deterministic_plan(
        card,
        options=RefinementOptions(),
        profile=UserProfileContext(planning_granularity="detailed"),
        evidence_summary=["报告必须包含数据分析和参考文献"],
    )

    assert len(plan.items) >= 5
    assert {item.kind for item in plan.items} >= {"milestone", "work_block", "step"}
    assert not validate_plan(plan, card)
    assert all(
        datetime.fromisoformat(value) <= deadline
        for item in plan.items
        for value in [item.start_time, item.deadline]
        if value
    )
    assert any(item.reminder_enabled for item in plan.items if item.kind == "milestone")


def test_rule_refinement_without_deadline_never_creates_absolute_reminder() -> None:
    card = _card()
    plan = deterministic_plan(
        card,
        options=RefinementOptions(),
        profile=None,
        evidence_summary=[],
    )

    assert plan.items
    assert all(item.start_time is None and item.deadline is None for item in plan.items)
    assert all(not item.reminder_enabled for item in plan.items)
    assert all(not item.need_confirm for item in plan.items)
    assert any("截止时间" in warning for warning in plan.warnings)


def test_profile_changes_granularity_but_not_parent_facts() -> None:
    deadline = (datetime.now(timezone.utc) + timedelta(days=3)).isoformat()
    card = _card(deadline=deadline)
    concise = deterministic_plan(
        card,
        options=RefinementOptions(),
        profile=UserProfileContext(planning_granularity="concise"),
        evidence_summary=[],
    )
    detailed = deterministic_plan(
        card,
        options=RefinementOptions(),
        profile=UserProfileContext(planning_granularity="detailed"),
        evidence_summary=[],
    )

    assert len(detailed.items) > len(concise.items)
    assert concise.parent_card_id == detailed.parent_card_id == card.id
    assert concise.objective == detailed.objective == card.title


def test_profile_changes_schedule_buffer_and_reminder_density_without_changing_facts() -> None:
    deadline = (datetime.now(timezone.utc) + timedelta(days=8)).isoformat()
    card = _card(deadline=deadline)
    light = deterministic_plan(
        card,
        options=RefinementOptions(),
        profile=UserProfileContext(
            active_period="morning",
            reminder_style="key_only",
            work_rhythm="steady",
            buffer_preference="generous",
            timezone="Asia/Shanghai",
        ),
        evidence_summary=["研究报告必须包含数据分析和参考文献"],
    )
    intensive = deterministic_plan(
        card,
        options=RefinementOptions(),
        profile=UserProfileContext(
            active_period="evening",
            reminder_style="multi_stage",
            work_rhythm="sprint",
            buffer_preference="compact",
            timezone="Asia/Shanghai",
        ),
        evidence_summary=["研究报告必须包含数据分析和参考文献"],
    )

    assert light.objective == intensive.objective == card.title
    assert light.parent_card_id == intensive.parent_card_id == card.id
    assert sum(item.reminder_enabled for item in light.items) == 1
    assert sum(item.reminder_enabled for item in intensive.items) >= 2
    assert light.profile_effects != intensive.profile_effects
    assert not light.constraint_errors
    assert not intensive.constraint_errors


def test_plan_validation_rejects_dependency_cycle_and_parent_fact_rewrite() -> None:
    card = _card(deadline=(datetime.now(timezone.utc) + timedelta(days=2)).isoformat())
    first = PlanItem(
        id="first",
        kind="step",
        title="第一步",
        dependencies=["second"],
    )
    second = PlanItem(
        id="second",
        kind="step",
        title="第二步",
        order=1,
        dependencies=["first"],
    )
    plan = CardRefinementPlan(
        id="cycle",
        parent_card_id=card.id,
        objective="被模型改写的标题",
        items=[first, second],
    )

    errors = validate_plan(plan, card)

    assert "计划包含循环依赖" in errors
    assert "计划修改了父卡标题事实" in errors


def test_text_and_markdown_extraction(tmp_path: Path) -> None:
    source = tmp_path / "requirements.md"
    source.write_text("# 作业要求\n截止周五，提交 PDF。", encoding="utf-8")

    result = asyncio.run(
        extract_document(
            source,
            name=source.name,
            declared_mime="text/markdown",
            attachment_id="attachment-md",
        )
    )

    assert result.descriptor.extraction_status == "succeeded"
    assert "截止周五" in result.text
    assert result.descriptor.sha256


def test_mime_spoof_is_rejected(tmp_path: Path) -> None:
    source = tmp_path / "fake.png"
    source.write_text("this is text, not an image", encoding="utf-8")

    result = asyncio.run(
        extract_document(
            source,
            name=source.name,
            declared_mime="image/png",
            attachment_id="attachment-spoof",
        )
    )

    assert result.descriptor.extraction_status == "unsupported"
    assert "不一致" in (result.descriptor.warning or "")


def test_docx_pptx_xlsx_and_pdf_extractors(tmp_path: Path) -> None:
    import fitz
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    docx_path = tmp_path / "requirements.docx"
    document = Document()
    document.add_paragraph("报告需要包含研究背景和参考文献")
    document.save(docx_path)

    pptx_path = tmp_path / "meeting.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "周三完成会议彩排"
    presentation.save(pptx_path)

    xlsx_path = tmp_path / "tasks.xlsx"
    workbook = Workbook()
    workbook.active.append(["事项", "截止"])
    workbook.active.append(["提交数据表", "周四 18:00"])
    workbook.save(xlsx_path)

    pdf_path = tmp_path / "brief.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Submit final report before Friday")
    pdf.save(pdf_path)
    pdf.close()

    async def extract_all():
        return await asyncio.gather(
            extract_document(
                docx_path,
                name=docx_path.name,
                declared_mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                attachment_id="docx",
            ),
            extract_document(
                pptx_path,
                name=pptx_path.name,
                declared_mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                attachment_id="pptx",
            ),
            extract_document(
                xlsx_path,
                name=xlsx_path.name,
                declared_mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                attachment_id="xlsx",
            ),
            extract_document(
                pdf_path,
                name=pdf_path.name,
                declared_mime="application/pdf",
                attachment_id="pdf",
            ),
        )

    docx_result, pptx_result, xlsx_result, pdf_result = asyncio.run(extract_all())
    assert "研究背景" in docx_result.text
    assert "会议彩排" in pptx_result.text
    assert "提交数据表" in xlsx_result.text
    assert "Submit final report" in pdf_result.text
    assert all(
        result.descriptor.extraction_status == "succeeded"
        for result in [docx_result, pptx_result, xlsx_result, pdf_result]
    )


def test_encrypted_pdf_is_rejected_without_leaking_content(tmp_path: Path) -> None:
    import fitz

    source = tmp_path / "protected.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "private assignment requirements")
    document.save(
        source,
        encryption=fitz.PDF_ENCRYPT_AES_256,
        owner_pw="owner-secret",
        user_pw="user-secret",
    )
    document.close()

    result = asyncio.run(
        extract_document(
            source,
            name=source.name,
            declared_mime="application/pdf",
            attachment_id="encrypted-pdf",
        )
    )

    assert result.descriptor.extraction_status == "password_protected"
    assert result.text == ""
    assert "加密" in (result.descriptor.warning or "")


def test_ooxml_path_traversal_is_rejected(tmp_path: Path) -> None:
    source = tmp_path / "unsafe.docx"
    with ZipFile(source, "w", compression=ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", "<document />")
        archive.writestr("../outside.txt", "unsafe")

    result = asyncio.run(
        extract_document(
            source,
            name=source.name,
            declared_mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            attachment_id="unsafe-docx",
        )
    )

    assert result.descriptor.extraction_status == "unsupported"
    assert result.text == ""
    assert "不安全路径" in (result.descriptor.warning or "")
